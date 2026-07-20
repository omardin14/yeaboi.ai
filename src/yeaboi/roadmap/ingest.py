"""Roadmap document ingestion — reads the quarterly roadmap from any source.

The Roadmap intake card lets the user point at wherever their quarterly roadmap
lives: a Confluence page, a Notion page, or a local file (.md/.txt/.rst/.pdf/
.docx/.pptx). This module turns a saved :class:`RoadmapSource` into plain text
for the analysis LLM call, degrading gracefully — every failure becomes a
warning string, never an exception, so the TUI page always renders.

Word/PowerPoint support is optional (``uv sync --extra docs``); the libraries
are lazy-imported the same way pymupdf is for PDFs in tools/codebase.py.
"""

import logging
import re
from dataclasses import dataclass
from pathlib import Path

logger = logging.getLogger(__name__)

# Character budget for the roadmap text fed to the analysis prompt — roughly 6k
# tokens of evidence: big enough for a quarter roadmap, small enough to leave
# room for the prompt itself. Larger docs are truncated with a warning.
_MAX_ROADMAP_CHARS = 24_000

# Local file types the ingester can read. PDF needs `uv sync --extra pdf`,
# Word/PowerPoint need `uv sync --extra docs`.
_LOCAL_EXTENSIONS = frozenset({".md", ".txt", ".rst", ".pdf", ".docx", ".pptx"})


@dataclass(frozen=True)
class RoadmapSource:
    """Where the quarterly roadmap lives — the saved config value-object.

    Persisted as the singleton roadmap_config row (roadmap/store.py) so the
    card can re-analyze the same source without asking again.
    """

    source_type: str = ""  # "confluence" | "notion" | "local"
    locator: str = ""  # page id (confluence/notion) or file path (local)
    label: str = ""  # display name: page title / file name


# ---------------------------------------------------------------------------
# Locator parsing — accept pasted URLs, bare ids, or (Confluence) page titles
# ---------------------------------------------------------------------------


def parse_confluence_locator(raw: str) -> str:
    """Normalize a pasted Confluence reference to a page id (or pass a title through).

    Handles the shapes people actually paste:
      - https://x.atlassian.net/wiki/spaces/KEY/pages/12345/Q3-Roadmap → "12345"
      - https://x.atlassian.net/pages/viewpage.action?pageId=12345    → "12345"
      - "12345"                                                        → "12345"
      - anything else is treated as a page title (needs CONFLUENCE_SPACE_KEY).
    """
    raw = raw.strip()
    m = re.search(r"/pages/(\d+)", raw)
    if m:
        return m.group(1)
    m = re.search(r"[?&]pageId=(\d+)", raw)
    if m:
        return m.group(1)
    return raw


def parse_notion_locator(raw: str) -> str:
    """Normalize a pasted Notion reference to a page id.

    Notion URLs end in a 32-hex id (often suffixed to a slug, sometimes with
    dashes): https://notion.so/ws/Q3-Roadmap-0123456789abcdef0123456789abcdef.
    Returns the id (dashed form preserved if given), else the input unchanged.
    """
    raw = raw.strip().split("?", 1)[0]
    m = re.search(r"([0-9a-fA-F]{8}-[0-9a-fA-F]{4}-[0-9a-fA-F]{4}-[0-9a-fA-F]{4}-[0-9a-fA-F]{12})\s*$", raw)
    if m:
        return m.group(1)
    m = re.search(r"([0-9a-fA-F]{32})\s*$", raw)
    if m:
        return m.group(1)
    return raw


# ---------------------------------------------------------------------------
# Local-file extraction (docx/pptx are the only NEW parsers in the codebase)
# ---------------------------------------------------------------------------


def _extract_docx(path: Path) -> str | None:
    """Extract text from a Word document via python-docx (paragraphs + tables).

    Returns None if python-docx is not installed or the file cannot be read —
    the caller turns None into an actionable warning.
    """
    try:
        import docx  # optional dependency: uv sync --extra docs
    except ImportError:
        logger.debug("python-docx not installed — cannot read %s", path)
        return None
    try:
        document = docx.Document(str(path))
        parts = [p.text for p in document.paragraphs if p.text.strip()]
        # Roadmap docs often carry the actual plan in tables — flatten each row
        # to a tab-joined line so the LLM sees the columns side by side.
        for table in document.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    parts.append("\t".join(cells))
        text = "\n".join(parts).strip()
        return text or None
    except Exception:
        logger.warning("Failed to read Word document: %s", path, exc_info=True)
        return None


def _extract_pptx(path: Path) -> str | None:
    """Extract text from a PowerPoint deck via python-pptx.

    Emits a "--- Slide N ---" marker per slide, all shape text frames, and the
    speaker notes (roadmap decks often carry the real detail in the notes).
    Returns None if python-pptx is not installed or the file cannot be read.
    """
    try:
        import pptx  # optional dependency: uv sync --extra docs
    except ImportError:
        logger.debug("python-pptx not installed — cannot read %s", path)
        return None
    try:
        presentation = pptx.Presentation(str(path))
        parts: list[str] = []
        for idx, slide in enumerate(presentation.slides, start=1):
            slide_parts: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    frame_text = shape.text_frame.text.strip()
                    if frame_text:
                        slide_parts.append(frame_text)
            if getattr(slide, "has_notes_slide", False):
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_parts.append(f"Notes: {notes}")
            if slide_parts:
                parts.append(f"--- Slide {idx} ---")
                parts.extend(slide_parts)
        text = "\n".join(parts).strip()
        return text or None
    except Exception:
        logger.warning("Failed to read PowerPoint deck: %s", path, exc_info=True)
        return None


def _read_local_file(path_str: str) -> tuple[str, str, list[str]]:
    """Read a local roadmap file by suffix. Returns (text, label, warnings)."""
    path = Path(path_str).expanduser()
    label = path.name
    if not path.exists() or not path.is_file():
        return "", label, [f"Roadmap file not found: {path}"]

    suffix = path.suffix.lower()
    if suffix not in _LOCAL_EXTENSIONS:
        return "", label, [f"Unsupported roadmap file type '{suffix}' — use .md, .txt, .rst, .pdf, .docx, or .pptx."]

    if suffix in (".md", ".txt", ".rst"):
        try:
            return path.read_text(encoding="utf-8", errors="replace"), label, []
        except OSError as e:
            logger.warning("Failed to read roadmap file %s: %s", path, e)
            return "", label, [f"Could not read {path.name}: {e}"]

    if suffix == ".pdf":
        # Reuse the codebase tool's lazy pymupdf extractor (returns None when
        # the optional dep is missing or the PDF is unreadable).
        from yeaboi.tools.codebase import _read_pdf

        text = _read_pdf(str(path))
        if text is None:
            return "", label, [f"Could not read {path.name} — PDF support requires: uv sync --extra pdf"]
        return text, label, []

    extractor = _extract_docx if suffix == ".docx" else _extract_pptx
    text = extractor(path)
    if text is None:
        kind = "Word" if suffix == ".docx" else "PowerPoint"
        return "", label, [f"Could not read {path.name} — {kind} support requires: uv sync --extra docs"]
    return text, label, []


# ---------------------------------------------------------------------------
# Source dispatch
# ---------------------------------------------------------------------------


def ingest_source(source: RoadmapSource) -> tuple[str, str, list[str]]:
    """Read the roadmap text from the configured source.

    Returns (text, resolved_label, warnings). Never raises — an unreadable or
    unconfigured source yields empty text plus a warning the engine folds into
    RoadmapAnalysis.warnings, so the results view can explain what went wrong.
    """
    logger.info("ingest_source: type=%r locator=%r", source.source_type, source.locator)
    warnings: list[str] = []

    if source.source_type == "confluence":
        # Lazy import — the atlassian SDK is an optional install path.
        from yeaboi.tools.confluence import confluence_read_page_text

        locator = parse_confluence_locator(source.locator)
        if locator.isdigit():
            result = confluence_read_page_text(page_id=locator)
        else:
            result = confluence_read_page_text(page_title=locator)
        text = result.get("text", "")
        label = result.get("title", "") or source.label or locator
        if result.get("error"):
            warnings.append(f"Confluence: {result['error']}")
    elif source.source_type == "notion":
        from yeaboi.tools.notion import notion_read_page_text

        result = notion_read_page_text(parse_notion_locator(source.locator))
        text = result.get("text", "")
        label = result.get("title", "") or source.label or source.locator
        if result.get("error"):
            warnings.append(f"Notion: {result['error']}")
    elif source.source_type == "local":
        text, label, warnings = _read_local_file(source.locator)
    else:
        return "", source.label, [f"Unknown roadmap source type: {source.source_type!r}"]

    if len(text) > _MAX_ROADMAP_CHARS:
        text = text[:_MAX_ROADMAP_CHARS]
        warnings.append(
            f"Roadmap truncated at {_MAX_ROADMAP_CHARS:,} characters — "
            "the analysis covers the first part of the document."
        )
    logger.info("ingest_source: %d chars from %r (%d warning(s))", len(text), label, len(warnings))
    return text, label, warnings
